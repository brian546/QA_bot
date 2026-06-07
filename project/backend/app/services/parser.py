from __future__ import annotations

import csv
import io
import os
from collections import Counter

from pypdf import PdfReader

SUPPORTED_UPLOAD_EXTENSIONS = {".pdf", ".txt", ".md", ".markdown", ".csv", ".docx", ".pptx"}


def _strip_repeated_page_edges(page_texts: list[str]) -> list[str]:
    """Remove repeated first and last lines that likely represent headers/footers."""
    first_lines: list[str] = []
    last_lines: list[str] = []

    split_pages: list[list[str]] = []
    for text in page_texts:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        split_pages.append(lines)
        if lines:
            first_lines.append(lines[0])
            last_lines.append(lines[-1])

    first_common = {line for line, count in Counter(first_lines).items() if count > 1}
    last_common = {line for line, count in Counter(last_lines).items() if count > 1}

    cleaned: list[str] = []
    for lines in split_pages:
        if lines and lines[0] in first_common:
            lines = lines[1:]
        if lines and lines[-1] in last_common:
            lines = lines[:-1]
        cleaned.append("\n".join(lines).strip())
    return cleaned


def parse_pdf_pages(pdf_bytes: bytes, filename: str) -> list[dict[str, object]]:
    """Parse PDF by page and keep source metadata."""
    reader = PdfReader(io.BytesIO(pdf_bytes))
    raw_pages: list[str] = []
    for page in reader.pages:
        raw_pages.append((page.extract_text() or "").strip())

    cleaned_pages = _strip_repeated_page_edges(raw_pages)

    parsed: list[dict[str, object]] = []
    for idx, text in enumerate(cleaned_pages, start=1):
        if not text:
            continue
        parsed.append({"filename": filename, "page": idx, "text": text})
    return parsed


def _decode_text(raw_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw_bytes.decode("utf-8", errors="replace")


def _parse_text_document(raw_bytes: bytes, filename: str) -> list[dict[str, object]]:
    text = _decode_text(raw_bytes).strip()
    if not text:
        return []
    return [{"filename": filename, "page": 1, "text": text}]


def _parse_csv_document(raw_bytes: bytes, filename: str) -> list[dict[str, object]]:
    decoded = _decode_text(raw_bytes)
    rows: list[str] = []
    reader = csv.reader(io.StringIO(decoded))
    for row in reader:
        cells = [cell.strip() for cell in row]
        if any(cells):
            rows.append(" | ".join(cells))

    text = "\n".join(rows).strip()
    if not text:
        return []
    return [{"filename": filename, "page": 1, "text": text}]


def _parse_docx_document(raw_bytes: bytes, filename: str) -> list[dict[str, object]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise ValueError("DOCX parsing requires python-docx.") from exc

    document = Document(io.BytesIO(raw_bytes))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    if not text:
        return []
    return [{"filename": filename, "page": 1, "text": text}]


def _parse_pptx_document(raw_bytes: bytes, filename: str) -> list[dict[str, object]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("PPTX parsing requires python-pptx.") from exc

    presentation = Presentation(io.BytesIO(raw_bytes))
    pages: list[dict[str, object]] = []

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        fragments: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = str(text).strip()
            if text:
                fragments.append(text)

        slide_text = "\n".join(fragments).strip()
        if slide_text:
            pages.append({"filename": filename, "page": slide_idx, "text": slide_text})

    return pages


def parse_document_pages(raw_bytes: bytes, filename: str) -> list[dict[str, object]]:
    """Parse supported document formats into page-like records."""
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        return parse_pdf_pages(raw_bytes, filename)
    if ext in {".txt", ".md", ".markdown"}:
        return _parse_text_document(raw_bytes, filename)
    if ext == ".csv":
        return _parse_csv_document(raw_bytes, filename)
    if ext == ".docx":
        return _parse_docx_document(raw_bytes, filename)
    if ext == ".pptx":
        return _parse_pptx_document(raw_bytes, filename)

    return []
